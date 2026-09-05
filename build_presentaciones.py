"""Genera las presentaciones (.pptx y .pdf) de los tres casos de uso.

Las figuras se extraen directamente de los notebooks ya ejecutados, de modo que
las presentaciones nunca se desincronizan del análisis: si se vuelve a ejecutar
un notebook y cambia un gráfico, basta con volver a lanzar este script.

Uso:
    conda run -n dataanalyst-final python build_presentaciones.py

La conversión a PDF usa PowerPoint mediante COM (solo Windows). Si no está
disponible, se generan los .pptx y se avisa para exportarlos a mano.
"""

import base64
import json
import os
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

RAIZ = Path(__file__).parent

AZUL_OSCURO = RGBColor(0x1F, 0x35, 0x54)
AZUL = RGBColor(0x4C, 0x72, 0xB0)
ROJO = RGBColor(0xC4, 0x4E, 0x52)
GRIS = RGBColor(0x55, 0x5D, 0x66)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)

ANCHO = Inches(13.333)
ALTO = Inches(7.5)


# --------------------------------------------------------------------------
# Extracción de figuras desde los notebooks
# --------------------------------------------------------------------------

def extraer_figuras(notebook, destino):
    """Guarda como PNG cada imagen de salida del notebook. Devuelve {celda: ruta}."""
    destino.mkdir(parents=True, exist_ok=True)
    nb = json.loads(Path(notebook).read_text(encoding='utf-8'))
    figuras = {}
    for i, celda in enumerate(nb['cells']):
        for salida in celda.get('outputs', []):
            png = salida.get('data', {}).get('image/png')
            if not png:
                continue
            ruta = destino / f'fig_{i:02d}.png'
            ruta.write_bytes(base64.b64decode(png))
            figuras[i] = ruta
    return figuras


# --------------------------------------------------------------------------
# Construcción de diapositivas
# --------------------------------------------------------------------------

def nueva_presentacion():
    prs = Presentation()
    prs.slide_width = ANCHO
    prs.slide_height = ALTO
    return prs


def _caja(slide, texto, izq, arr, ancho, alto, tam=18, negrita=False,
          color=GRIS, alineacion=PP_ALIGN.LEFT, interlineado=1.25,
          espacio_despues=0, sangria=None, anclaje=None):
    caja = slide.shapes.add_textbox(izq, arr, ancho, alto)
    marco = caja.text_frame
    marco.word_wrap = True
    if anclaje is not None:
        marco.vertical_anchor = anclaje
    for n, linea in enumerate(texto.split('\n')):
        p = marco.paragraphs[0] if n == 0 else marco.add_paragraph()
        p.text = linea
        p.alignment = alineacion
        p.line_spacing = interlineado
        if espacio_despues:
            p.space_after = Pt(espacio_despues)
        if sangria is not None:
            # Sangría francesa: la primera línea sale al margen y las líneas
            # continuadas se alinean bajo el texto, no bajo la viñeta.
            pPr = p._p.get_or_add_pPr()
            pPr.set('marL', str(int(sangria)))
            pPr.set('indent', str(int(-sangria)))
        for run in p.runs:
            run.font.size = Pt(tam)
            run.font.bold = negrita
            run.font.color.rgb = color
    return caja


def portada(prs, titulo, subtitulo, pie):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fondo = slide.shapes.add_shape(1, 0, 0, ANCHO, ALTO)
    fondo.fill.solid()
    fondo.fill.fore_color.rgb = AZUL_OSCURO
    fondo.line.fill.background()

    _caja(slide, titulo, Inches(1), Inches(2.4), Inches(11.3), Inches(1.6),
          tam=40, negrita=True, color=BLANCO)
    _caja(slide, subtitulo, Inches(1), Inches(4.0), Inches(11.3), Inches(1.2),
          tam=20, color=RGBColor(0xB8, 0xC6, 0xD9))
    _caja(slide, pie, Inches(1), Inches(6.4), Inches(11.3), Inches(0.5),
          tam=13, color=RGBColor(0x8A, 0x9B, 0xB3))
    return slide


def seccion(prs, numero, titulo):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    barra = slide.shapes.add_shape(1, 0, Inches(3.1), ANCHO, Inches(1.3))
    barra.fill.solid()
    barra.fill.fore_color.rgb = AZUL_OSCURO
    barra.line.fill.background()
    _caja(slide, f'{numero}   {titulo}', Inches(1), Inches(3.35), Inches(11.3),
          Inches(1), tam=28, negrita=True, color=BLANCO)
    return slide


def _encabezado(slide, titulo):
    _caja(slide, titulo, Inches(0.7), Inches(0.4), Inches(12), Inches(0.8),
          tam=26, negrita=True, color=AZUL_OSCURO)
    linea = slide.shapes.add_shape(1, Inches(0.7), Inches(1.15), Inches(1.8), Emu(28000))
    linea.fill.solid()
    linea.fill.fore_color.rgb = ROJO
    linea.line.fill.background()


def grafico(prs, titulo, imagen, comentario=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _encabezado(slide, titulo)

    arriba = Inches(1.45)
    # El pie ocupa la franja inferior; la imagen se escala para no invadirla.
    alto_max = Inches(4.55) if comentario else Inches(5.6)
    ancho_max = Inches(10.9)

    with Image.open(imagen) as img:
        px_ancho, px_alto = img.size
    escala = min(ancho_max / px_ancho, alto_max / px_alto)
    ancho = int(px_ancho * escala)
    alto = int(px_alto * escala)

    slide.shapes.add_picture(str(imagen), Emu(int((ANCHO - ancho) / 2)), arriba,
                             width=Emu(ancho), height=Emu(alto))

    if comentario:
        _caja(slide, comentario, Inches(0.9), Inches(6.25), Inches(11.6), Inches(0.95),
              tam=15, color=GRIS, interlineado=1.2)
    return slide


def puntos(prs, titulo, lineas, destacado=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _encabezado(slide, titulo)
    arriba = Inches(1.6)
    if destacado:
        caja = slide.shapes.add_shape(1, Inches(0.7), arriba, Inches(12), Inches(1.0))
        caja.fill.solid()
        caja.fill.fore_color.rgb = RGBColor(0xF2, 0xE4, 0xE4)
        caja.line.color.rgb = ROJO
        _caja(slide, destacado, Inches(1.0), arriba + Inches(0.18), Inches(11.4),
              Inches(0.7), tam=18, negrita=True, color=ROJO)
        arriba = arriba + Inches(1.35)

    texto = '\n'.join(f'•  {l}' for l in lineas)
    alto = Emu(int(ALTO - arriba - Inches(0.8)))
    _caja(slide, texto, Inches(0.9), arriba, Inches(11.6), alto,
          tam=17, color=GRIS, interlineado=1.2, espacio_despues=16,
          sangria=Inches(0.35), anclaje=MSO_ANCHOR.MIDDLE)
    return slide


def cifras(prs, titulo, tarjetas):
    """tarjetas: lista de (valor, etiqueta). Máximo 4."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _encabezado(slide, titulo)
    ancho_t = Inches(2.9)
    hueco = Inches(0.32)
    total = len(tarjetas) * ancho_t + (len(tarjetas) - 1) * hueco
    izq = Emu(int((ANCHO - total) / 2))
    for valor, etiqueta in tarjetas:
        caja = slide.shapes.add_shape(1, izq, Inches(2.3), ancho_t, Inches(2.5))
        caja.fill.solid()
        caja.fill.fore_color.rgb = RGBColor(0xF4, 0xF6, 0xF9)
        caja.line.color.rgb = RGBColor(0xD8, 0xDE, 0xE8)
        _caja(slide, valor, izq, Inches(2.75), ancho_t, Inches(1.0),
              tam=40, negrita=True, color=AZUL, alineacion=PP_ALIGN.CENTER)
        _caja(slide, etiqueta, izq + Inches(0.2), Inches(3.85), ancho_t - Inches(0.4),
              Inches(0.9), tam=13, color=GRIS, alineacion=PP_ALIGN.CENTER)
        izq = Emu(int(izq + ancho_t + hueco))
    return slide


def cierre(prs, titulo, mensaje, acciones):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _encabezado(slide, titulo)
    caja = slide.shapes.add_shape(1, Inches(0.7), Inches(1.6), Inches(12), Inches(1.3))
    caja.fill.solid()
    caja.fill.fore_color.rgb = AZUL_OSCURO
    caja.line.fill.background()
    _caja(slide, mensaje, Inches(1.0), Inches(1.85), Inches(11.4), Inches(1.0),
          tam=19, negrita=True, color=BLANCO)
    texto = '\n'.join(f'{n}.  {a}' for n, a in enumerate(acciones, 1))
    _caja(slide, texto, Inches(0.9), Inches(3.25), Inches(11.6), Inches(3.5),
          tam=17, color=GRIS, interlineado=1.2, espacio_despues=14,
          sangria=Inches(0.42), anclaje=MSO_ANCHOR.MIDDLE)
    return slide


# --------------------------------------------------------------------------
# Presentación 1 — Caso principal
# --------------------------------------------------------------------------

def presentacion_principal():
    figs = extraer_figuras(RAIZ / 'main_project/02_analisis_operadores.ipynb',
                           RAIZ / 'main_project/presentacion/figuras')
    prs = nueva_presentacion()

    portada(prs, 'CallMeMaybe: ¿quiénes son los operadores ineficaces?',
            'Análisis del servicio de telefonía virtual  ·  agosto – noviembre 2019',
            'Proyecto final de analítica de datos  ·  Caso principal')

    puntos(prs, 'El encargo',
           ['CallMeMaybe quiere dar a los supervisores una función que señale a sus '
            'operadores menos eficaces.',
            'Un operador sería ineficaz por tres señales: muchas llamadas entrantes '
            'perdidas, esperas largas y pocas llamadas salientes.',
            'La pregunta de partida: ¿pueden medirse esas tres señales con los datos '
            'disponibles?'])

    cifras(prs, 'El punto de partida',
           [('806 709', 'llamadas analizadas'),
            ('307', 'organizaciones clientes'),
            ('1 092', 'operadores'),
            ('4 meses', 'de actividad registrada')])

    seccion(prs, '01', 'El hallazgo que cambió el análisis')

    grafico(prs, 'Más de la mitad de las llamadas entrantes se pierden', figs[7],
            'Composición del tráfico. En las llamadas salientes, "perdida" significa que '
            'el destinatario no respondió, algo ajeno al operador.')

    puntos(prs, 'Las llamadas perdidas no tienen dueño',
           ['De las 104 323 llamadas entrantes perdidas, 103 397 no tienen ningún '
            'operador asignado.',
            'Es coherente: una llamada perdida es, por definición, una llamada que '
            'nadie atendió.',
            'Solo 926 llamadas perdidas son atribuibles, repartidas entre 239 '
            'operadores: una base demasiado escasa para un ranking.',
            'Afecta a 305 de los 307 clientes, así que no es un problema de unas pocas '
            'organizaciones.'],
           destacado='El 99% de las llamadas perdidas no puede atribuirse a ninguna persona')

    puntos(prs, 'La consecuencia: dos niveles de análisis',
           ['NIVEL ORGANIZACIÓN — la pérdida de llamadas se analiza por cliente. Es un '
            'problema de cobertura, no de desempeño individual.',
            'NIVEL OPERADOR — la evaluación personal se sostiene sobre lo que sí es '
            'atribuible: la espera en las llamadas que sí atiende y su volumen de '
            'llamadas salientes.',
            'Construir el ranking sobre las llamadas perdidas habría producido un '
            'resultado que parece riguroso y no lo es.'])

    seccion(prs, '02', 'Qué operadores están fallando')

    grafico(prs, 'Cada operador se compara solo con sus pares', figs[22],
            'La distribución es claramente bimodal: hay operadores que solo reciben '
            'llamadas y otros que solo emiten. Comparar entre grupos distintos sería un error.')

    grafico(prs, 'Los señalados son un grupo claramente distinto', figs[32],
            'Esperan casi el doble por llamada y gestionan una fracción del volumen '
            'diario del resto, con la misma antigüedad en el servicio.')

    cifras(prs, 'El resultado',
           [('91', 'operadores señalados'),
            ('12,6%', 'del total evaluado'),
            ('49', 'clientes afectados'),
            ('1,2%', 'del volumen que gestionan')])

    seccion(prs, '03', 'Lo que confirman las pruebas estadísticas')

    grafico(prs, 'Quien más atiende es quien más hace esperar', figs[41],
            'Kruskal-Wallis p < 0,001. Los tres perfiles difieren entre sí tras corregir '
            'por comparaciones múltiples. Apunta a saturación, no a falta de aptitud.')

    grafico(prs, 'El plan contratado condiciona la calidad del servicio', figs[44],
            'Kruskal-Wallis p < 0,001. Los clientes del plan B pierden significativamente '
            'más llamadas entrantes que los del plan C.')

    cierre(prs, 'Conclusiones y recomendaciones',
           'La función pedida es viable, pero no resuelve el problema principal del servicio',
           ['Registrar el operador —o la cola— también en las llamadas perdidas. Sin ese '
            'dato, la señal más importante seguirá siendo inmedible.',
            'Añadir al panel del supervisor un indicador de cobertura por cliente, junto '
            'al ranking individual.',
            'Revisar las condiciones del plan B, cuyos clientes pierden más llamadas.',
            'Usar el ranking como lista corta de revisión, no como evaluación de desempeño.',
            'Revisar el reparto de la carga de llamadas entrantes antes que el desempeño '
            'individual.'])

    puntos(prs, 'Limitaciones del análisis',
           ['El período cubre cuatro meses de 2019 y coincide con una fase de crecimiento '
            'de la base de clientes.',
            'No se conoce el rol asignado a cada operador: el perfil se infiere del '
            'comportamiento observado.',
            'El índice mide desempeño relativo dentro de cada perfil. Si un equipo entero '
            'funcionara mal, solo se señalaría a los peores de ese equipo.',
            'Los operadores salientes sin llamadas entrantes solo pueden evaluarse en una '
            'dimensión.'])

    return prs, RAIZ / 'main_project/presentacion/CallMeMaybe_operadores_ineficaces.pptx'


# --------------------------------------------------------------------------
# Presentación 2 — Test A/B
# --------------------------------------------------------------------------

def presentacion_ab():
    figs = extraer_figuras(RAIZ / 'ab_project/ab_test_analysis.ipynb',
                           RAIZ / 'ab_project/presentacion/figuras')
    prs = nueva_presentacion()

    portada(prs, '¿Funciona el nuevo sistema de recomendaciones?',
            'Evaluación de la prueba A/B recommender_system_test  ·  dic 2020 – ene 2021',
            'Proyecto final de analítica de datos  ·  Caso de test A/B')

    puntos(prs, 'La situación',
           ['Se lanzó una prueba A/B para validar un sistema de recomendaciones mejorado.',
            'El equipo que la puso en marcha abandonó el proyecto y solo dejó la '
            'especificación técnica y los datos.',
            'Se esperaba una mejora de al menos el 10% en cada etapa del embudo de compra.',
            'La pregunta: ¿qué se puede concluir de esa prueba?'])

    seccion(prs, '01', 'Primero: ¿es fiable la prueba?')

    puntos(prs, 'La prueba incumple su propia especificación',
           ['MUESTRA — 3 675 participantes frente a los 6 000 previstos.',
            'REPARTO — el grupo A casi triplica al B, incompatible con una asignación '
            'aleatoria equilibrada.',
            'AUDIENCIA — 194 participantes fuera de la región EU, cuando la prueba debía '
            'limitarse a ella.',
            'AISLAMIENTO — 887 participantes están también en una segunda prueba '
            'simultánea.',
            'OBSERVACIÓN — los datos terminan el 30 de diciembre; solo el 55% de los '
            'usuarios completa los 14 días exigidos.'],
           destacado='Cinco de los seis parámetros declarados no se cumplen')

    grafico(prs, 'Además, una campaña navideña se solapa con la prueba', figs[10],
            'La promoción de Navidad afecta al comportamiento de compra justo en el tramo '
            'final de medición. El 25 de diciembre no hay ningún evento registrado.')

    seccion(prs, '02', 'Qué dicen los resultados')

    grafico(prs, 'El grupo B no supera al A en ninguna etapa', figs[19],
            'Frente al +10% esperado, todas las diferencias son negativas. El retroceso en '
            'product_page es estadísticamente significativo (p < 0,001).')

    puntos(prs, 'Pero hay algo más importante que el resultado',
           ['Para detectar una mejora del 10% en el carrito o en la compra harían falta '
            'unos 3 500 usuarios por grupo.',
            'El grupo B tuvo 655.',
            'La potencia estadística alcanzada en esas etapas ronda el 30%, frente al 80% '
            'que se considera estándar.',
            'El tamaño previsto en la propia especificación —6 000 participantes— tampoco '
            'habría bastado.'],
           destacado='Aunque el sistema hubiera funcionado, la prueba lo habría pasado por '
                     'alto 7 de cada 10 veces')

    cifras(prs, 'La prueba en cifras',
           [('3:1', 'desequilibrio A vs B'),
            ('−14%', 'peor resultado de B'),
            ('30%', 'potencia estadística'),
            ('5 de 6', 'parámetros incumplidos')])

    cierre(prs, 'Recomendación',
           'No adoptar el sistema, no descartarlo tampoco: repetir la prueba correctamente',
           ['Reclutar al menos 3 500 usuarios por grupo, el mínimo para detectar el efecto '
            'buscado.',
            'Verificar el equilibrio de la asignación antes de lanzar.',
            'Excluir a los usuarios inscritos en otras pruebas simultáneas.',
            'Filtrar la audiencia por región en el momento de asignar, no después.',
            'Mantener la recogida de datos 14 días más allá del cierre de admisiones.',
            'Evitar el solapamiento con campañas de marketing y auditar el registro del '
            'evento de carrito.'])

    puntos(prs, 'Por qué no basta con decir "no funciona"',
           ['El único resultado sólido —el retroceso en product_page— justifica no '
            'desplegar el cambio tal como está.',
            'Pero en las etapas que de verdad importan para el negocio, la prueba nunca '
            'tuvo capacidad de detectar nada.',
            'Declarar que el sistema no funciona sería un error simétrico al de adoptarlo: '
            'ambas decisiones se apoyarían en evidencia que no las sostiene.'])

    return prs, RAIZ / 'ab_project/presentacion/Test_AB_sistema_recomendaciones.pptx'


# --------------------------------------------------------------------------
# Presentación 3 — Caso SQL
# --------------------------------------------------------------------------

def presentacion_sql():
    figs = extraer_figuras(RAIZ / 'sql_project/sql_analysis.ipynb',
                           RAIZ / 'sql_project/presentacion/figuras')
    prs = nueva_presentacion()

    portada(prs, 'Un servicio para lectores: ¿dónde está la oportunidad?',
            'Análisis de la base de datos de libros, autores, calificaciones y reseñas',
            'Proyecto final de analítica de datos  ·  Caso SQL')

    puntos(prs, 'El encargo',
           ['La pandemia desplazó tiempo de ocio hacia la lectura y multiplicó las '
            'aplicaciones para lectores.',
            'Se dispone de la base de datos de uno de los servicios que compiten en ese '
            'mercado.',
            'El objetivo: fundamentar la propuesta de valor de un nuevo producto a partir '
            'de lo que los datos revelan.'])

    cifras(prs, 'La base de datos',
           [('1 000', 'libros'),
            ('636', 'autores'),
            ('6 456', 'calificaciones'),
            ('2 793', 'reseñas escritas')])

    seccion(prs, '01', 'Qué contiene el catálogo')

    grafico(prs, 'Un catálogo contemporáneo, no un fondo histórico', figs[12],
            '819 de los 1 000 libros —el 82%— se publicaron después del año 2000, con el '
            'grueso concentrado en esa década.')

    grafico(prs, 'Pocas reseñas y calificaciones uniformemente altas', figs[16],
            'Ningún libro supera las siete reseñas. Las calificaciones se agolpan entre 4 '
            'y 5, de modo que apenas distinguen unos títulos de otros.')

    grafico(prs, 'La oferta editorial está concentrada', figs[20],
            'Penguin Books encabeza el catálogo y aparece además con su sello Penguin '
            'Classics: es el interlocutor prioritario, pero también una dependencia.')

    seccion(prs, '02', 'Quién genera el contenido')

    puntos(prs, 'Un núcleo mínimo sostiene la comunidad',
           ['Solo 6 usuarios han calificado más de 50 libros.',
            'Ese grupo escribe 24,3 reseñas de media, muy por encima del resto de la base.',
            'Toda reseña lleva asociada una calificación, pero no al revés: escribir es el '
            'gesto costoso.',
            'Solo 19 libros de 1 000 superan las 50 calificaciones.'],
           destacado='El contenido que da valor al servicio lo produce un puñado de personas')

    cierre(prs, 'Propuesta de valor',
           'El problema del mercado no es la falta de libros, sino la falta de criterio '
           'para elegirlos',
           ['Convertir la reseña en el eje del producto, no en un accesorio: pedirla al '
            'terminar un libro, en formatos breves y guiados.',
            'Cultivar y dar visibilidad al núcleo de usuarios prolíficos: su pérdida '
            'vaciaría el contenido del servicio.',
            'Recomendar por afinidad entre lectores, no por nota media, que con estas '
            'calificaciones no distingue casi nada.',
            'Vigilar la dependencia de un único grupo editorial en la negociación de '
            'contenidos.'])

    puntos(prs, 'Limitaciones',
           ['La base contiene 1 000 libros y una comunidad reducida: es una muestra del '
            'mercado, no el mercado completo.',
            'No hay marcas temporales en calificaciones ni reseñas, de modo que no puede '
            'analizarse la evolución de la actividad.',
            'No se dispone de datos de lectura ni de compra, solo de valoración.'])

    return prs, RAIZ / 'sql_project/presentacion/Servicio_libros_propuesta_valor.pptx'


# --------------------------------------------------------------------------
# Conversión a PDF
# --------------------------------------------------------------------------

def a_pdf(rutas_pptx):
    """Convierte los .pptx a .pdf usando PowerPoint (Windows)."""
    try:
        import win32com.client
    except ImportError:
        print('\n[aviso] pywin32 no está instalado: no se pueden generar los PDF '
              'automáticamente.\n        Abre cada .pptx y usa Archivo → Exportar → PDF.')
        return []

    generados = []
    app = win32com.client.Dispatch('PowerPoint.Application')
    try:
        for ruta in rutas_pptx:
            pdf = ruta.with_suffix('.pdf')
            deck = app.Presentations.Open(str(ruta.resolve()), WithWindow=False)
            deck.SaveAs(str(pdf.resolve()), 32)  # 32 = formato PDF
            deck.Close()
            generados.append(pdf)
            print(f'  PDF generado: {pdf.name}')
    finally:
        app.Quit()
    return generados


if __name__ == '__main__':
    rutas = []
    for constructor in (presentacion_principal, presentacion_ab, presentacion_sql):
        prs, ruta = constructor()
        ruta.parent.mkdir(parents=True, exist_ok=True)
        prs.save(ruta)
        print(f'{ruta.relative_to(RAIZ)}  ({len(prs.slides)} diapositivas)')
        rutas.append(ruta)

    print('\nConvirtiendo a PDF...')
    a_pdf(rutas)
    print('\nListo.')
